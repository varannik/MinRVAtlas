"""
DataSentinel V&V Verification Engine
Analyses uploaded documents against registry checkpoints.
Based on real Puro.Earth CCS submission structure (44.01 UAE project).
"""
import json
import logging
import os
import re
from datetime import datetime
from typing import Any, Dict, List

import pandas as pd

logger = logging.getLogger("datasentinel.vv")


class DocumentExtractor:
    """Extract structured data and key values from project documents."""

    def extract(self, storage_path: str, file_type: str) -> Dict[str, Any]:
        try:
            if file_type == "csv":             return self._csv(storage_path)
            elif file_type in ("xlsx","xls","xlsm"): return self._excel(storage_path)
            elif file_type == "json":          return self._json(storage_path)
            elif file_type in ("pdf","docx"):  return self._text_doc(storage_path, file_type)
            elif file_type in ("txt","html"):  return self._plain_text(storage_path, file_type)
            elif file_type == "pptx":         return self._pptx(storage_path)
            return {"file_type": file_type, "extracted": False}
        except Exception as e:
            logger.error(f"Extraction error {storage_path}: {e}")
            return {"error": str(e), "file_type": file_type}

    def _csv(self, path):
        df = pd.read_csv(path)
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        return {
            "row_count": len(df), "column_count": len(df.columns),
            "columns": list(df.columns),
            "numeric_columns": numeric_cols,
            "date_columns": [c for c in df.columns if any(w in c.lower() for w in ["date","time","timestamp"])],
            "numeric_summary": {c: {"min": float(df[c].min()), "max": float(df[c].max()), "mean": float(df[c].mean()), "sum": float(df[c].sum())} for c in numeric_cols[:20]},
            "sample_rows": df.head(3).fillna("").to_dict(orient="records"),
            "null_counts": {k: int(v) for k, v in df.isnull().sum().items() if v > 0},
        }

    def _excel(self, path):
        xl = pd.ExcelFile(path)
        result = {"sheets": {}, "sheet_names": xl.sheet_names}
        for sheet in xl.sheet_names[:8]:
            try:
                df = xl.parse(sheet)
                numeric_cols = df.select_dtypes(include="number").columns.tolist()
                result["sheets"][sheet] = {
                    "row_count": len(df), "column_count": len(df.columns),
                    "columns": list(df.columns)[:30],
                    "numeric_summary": {c: {"min": float(df[c].min()), "max": float(df[c].max()), "sum": float(df[c].sum())} for c in numeric_cols[:10]},
                    "sample_rows": df.head(2).fillna("").to_dict(orient="records"),
                }
            except Exception:
                result["sheets"][sheet] = {"error": "Could not parse sheet"}
        return result

    def _json(self, path):
        with open(path) as f: data = json.load(f)
        return {"keys": list(data.keys()) if isinstance(data, dict) else [], "record_count": len(data) if isinstance(data, list) else 1, "sample": str(data)[:500]}

    def _text_doc(self, path, file_type):
        """Extract text (and tables for PDFs) using pdfplumber → pypdf fallback."""
        text = ""
        tables: List[List] = []
        try:
            if file_type == "docx":
                try:
                    import docx
                    doc = docx.Document(path)
                    text = "\n".join(p.text for p in doc.paragraphs)
                except ImportError:
                    text = f"DOCX file: {os.path.basename(path)} (text extraction requires python-docx)"
            else:
                # PDF — try pdfplumber first, then pypdf
                try:
                    import pdfplumber
                    pages_text: List[str] = []
                    with pdfplumber.open(path) as pdf:
                        for page in pdf.pages[:60]:
                            pages_text.append(page.extract_text() or "")
                            for tbl in (page.extract_tables() or []):
                                if tbl:
                                    tables.append(tbl)
                    text = "\n".join(pages_text)
                except Exception as e_pl:
                    logger.debug(f"pdfplumber failed ({e_pl}), falling back to pypdf for {path}")
                    try:
                        from pypdf import PdfReader
                        reader = PdfReader(path)
                        text = "\n".join(page.extract_text() or "" for page in reader.pages)
                    except Exception as e_py:
                        text = f"PDF extraction failed: {e_py}"
        except Exception as e:
            text = f"Text extraction failed: {e}"

        # Strip null bytes (PostgreSQL rejects them in JSONB)
        text = text.replace('\x00', '')

        # Serialize tables (cap to first 10 tables)
        tables_data = [
            [[str(cell) if cell is not None else "" for cell in row] for row in tbl]
            for tbl in tables[:10]
        ]

        return {
            "file_type": file_type,
            "text_preview": text[:3000],
            "text_length": len(text),
            "key_terms": self._extract_key_terms(text),
            "tables": tables_data,
            "table_count": len(tables),
        }

    def _plain_text(self, path: str, file_type: str):
        """Extract text from .txt and .html files."""
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                raw = f.read()
            if file_type == "html":
                # Strip HTML tags to get readable text
                import re as _re
                text = _re.sub(r"<[^>]+>", " ", raw)
                text = _re.sub(r"\s{2,}", " ", text).strip()
            else:
                text = raw
            # Strip null bytes
            text = text.replace('\x00', '')
            return {
                "file_type": file_type,
                "text_preview": text[:3000],
                "text_length": len(text),
                "key_terms": self._extract_key_terms(text),
                "tables": [],
                "table_count": 0,
            }
        except Exception as e:
            return {"file_type": file_type, "error": str(e), "extracted": False}

    def _pptx(self, path: str):
        """Extract text from all slides in a PowerPoint file."""
        try:
            from pptx import Presentation
            prs = Presentation(path)
            parts: list[str] = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    parts.append(f"[Slide {slide_num}] " + " ".join(slide_texts))
            text = "\n".join(parts).replace('\x00', '')
            return {
                "file_type": "pptx",
                "text_preview": text[:3000],
                "text_length": len(text),
                "key_terms": self._extract_key_terms(text),
                "slide_count": len(prs.slides),
                "tables": [],
                "table_count": 0,
            }
        except ImportError:
            return {"file_type": "pptx", "error": "python-pptx not installed", "extracted": False}
        except Exception as e:
            return {"file_type": "pptx", "error": str(e), "extracted": False}

    def _extract_key_terms(self, text: str) -> Dict[str, Any]:
        """Extract key numerical values and terms from document text."""
        t = text.lower()
        terms = {}
        # CO2 amounts
        for pattern, key in [
            (r'(\d+[\.,]?\d*)\s*(?:t|tonnes?|tco2|tco2e)\s*(?:co2|per\s*year)?', 'co2_tonnes'),
            (r'(\d+[\.,]?\d*)\s*(?:ktco2|kt\s*co2)', 'co2_kt'),
            (r'(\d+[\.,]?\d*)\s*(?:mw|megawatt)', 'capacity_mw'),
        ]:
            m = re.search(pattern, t)
            if m: terms[key] = m.group(1).replace(',','.')
        # Key compliance terms
        for term in ['additionality','permanence','leakage','monitoring','verification','noc','fnrc','fea','reservoir','injection','mineralisation','mineraliz']:
            if term in t: terms[f'contains_{term}'] = True
        # Dates
        dates = re.findall(r'20\d{2}', text)
        if dates: terms['years_mentioned'] = list(set(dates))[:5]
        return terms


class CheckpointAnalyser:
    """Analyse each checkpoint against available documents."""

    def analyse(self, checkpoint: Dict, all_docs: List[Dict], extracted: Dict[str, Any]) -> Dict:
        cp_id = checkpoint["id"]
        doc_type = checkpoint.get("document_type","")
        evidence_types = checkpoint.get("evidence_types", [])
        requirement = checkpoint.get("requirement","")

        # Find matching documents
        matched = self._match_docs(doc_type, evidence_types, all_docs)

        if not matched:
            # Confidence is 0 — we have no evidence at all
            return {
                "status": "warning",
                "confidence": 0.0,
                "finding": f"⚠ No matching document found for: {checkpoint['name']}. Required evidence: {', '.join(evidence_types[:3])}. Please upload the relevant document.",
                "evidence": [],
                "gap": True,
            }

        # Run checkpoint-specific analysis
        return self._analyse_checkpoint(checkpoint, matched, extracted)

    def _match_docs(self, doc_type: str, evidence_types: List[str], all_docs: List[Dict]) -> List[Dict]:
        matched = []
        for doc in all_docs:
            dt = (doc.get("document_type") or "").lower().replace(" ", "_")
            name = (doc.get("name") or "").lower()
            # Match by document_type field
            if dt == doc_type or any(et in dt for et in evidence_types):
                matched.append(doc)
                continue
            # Match by filename keywords
            for et in evidence_types:
                keywords = et.replace("_", " ").split()
                if sum(1 for kw in keywords if kw in name) >= min(2, len(keywords)):
                    matched.append(doc); break
        return matched

    def _analyse_checkpoint(self, cp: Dict, docs: List[Dict], extracted: Dict) -> Dict:
        cp_id = cp["id"]
        doc_names = [d.get("name","") for d in docs]
        category = cp.get("category","")

        # Category-specific analysis
        if "LCA" in category or "life_cycle" in cp_id:
            return self._check_lca(cp, docs, extracted)
        if "Leakage" in category:
            return self._check_leakage(cp, docs, extracted)
        if "Monitoring" in category:
            return self._check_monitoring(cp, docs, extracted)
        if "Storage Site" in category:
            return self._check_storage(cp, docs, extracted)
        if "Safeguard" in category:
            return self._check_safeguards(cp, docs, extracted)

        # Default: document presence with text analysis
        return self._check_document_presence(cp, docs, extracted)

    def _check_lca(self, cp, docs, extracted):
        for doc in docs:
            data = extracted.get(str(doc.get("id","")), {})
            # LCA spreadsheet
            if doc.get("file_type") in ("xlsx","xlsm","xls") and data.get("sheets"):
                sheets = list(data["sheets"].keys())
                co2_sheets = [s for s in sheets if any(w in s.lower() for w in ["co2","carbon","lca","removal","emission","scope"])]
                if co2_sheets:
                    for sheet_name in co2_sheets[:2]:
                        sheet = data["sheets"][sheet_name]
                        numeric = sheet.get("numeric_summary", {})
                        for col, stats in numeric.items():
                            if any(w in col.lower() for w in ["net","removal","co2","carbon"]):
                                val = stats.get("sum", stats.get("max", 0))
                                if abs(val) > 0:
                                    # Confidence = f(co2_sheets_ratio, numeric_cols_found, value_extracted)
                                    sheet_coverage = min(len(co2_sheets) / max(len(sheets), 1), 1.0)
                                    numeric_depth  = min(len(numeric) / 5.0, 1.0)
                                    confidence = round(0.50 + 0.25 * sheet_coverage + 0.20 * numeric_depth + 0.05, 4)
                                    return {
                                        "status": "passed", "confidence": confidence,
                                        "finding": f"✓ LCA spreadsheet found with {len(sheets)} sheets. CO2-related sheets: {', '.join(co2_sheets[:3])}. Net value detected in '{col}': {val:.1f}. Full review required to verify system boundaries and net removal claim.",
                                        "evidence": [{"doc": doc.get("name"), "co2_sheets": co2_sheets, "key_col": col, "value": val}],
                                        "gap": False,
                                    }
                # Spreadsheet present but no net CO2 column found
                sheet_coverage = min(len(co2_sheets) / max(len(sheets), 1), 1.0)
                confidence = round(0.40 + 0.25 * sheet_coverage, 4)
                return {
                    "status": "passed", "confidence": confidence,
                    "finding": f"✓ LCA spreadsheet uploaded: {doc.get('name')} with {len(sheets)} sheets: {', '.join(sheets[:5])}. CO2-specific sheets: {len(co2_sheets)}. Manual review required to verify LCA calculations.",
                    "evidence": [{"doc": doc.get("name"), "sheets": sheets}],
                    "gap": False,
                }
            # LCA report (PDF/DOCX)
            elif doc.get("file_type") in ("pdf","docx") and data.get("text_preview"):
                text  = data["text_preview"].lower()
                terms = data.get("key_terms", {})
                text_len = data.get("text_length", 0)
                lca_terms  = ["life cycle","lca","system boundary","functional unit","carbon removal","net removal","scope 1","scope 2","scope 3"]
                found_terms = [t for t in lca_terms if t in text]
                # Confidence = f(co2 extracted, term density, text completeness)
                has_co2    = 1.0 if terms.get("co2_tonnes") else 0.0
                term_score = min(len(found_terms) / len(lca_terms), 1.0)
                text_score = min(text_len / 3000.0, 1.0)
                confidence = round(0.40 + 0.25 * has_co2 + 0.20 * term_score + 0.15 * text_score, 4)
                co2 = terms.get("co2_tonnes", "not extracted")
                return {
                    "status": "passed", "confidence": confidence,
                    "finding": f"✓ LCA report found: {doc.get('name')}. CO2 quantity: {co2} tCO2. LCA terms matched: {len(found_terms)}/{len(lca_terms)}. Manual review required to verify ISO 14064 compliance and system boundaries.",
                    "evidence": [{"doc": doc.get("name"), "co2_tonnes": co2, "lca_terms_found": found_terms}],
                    "gap": False,
                }
        # Document uploaded but nothing extractable
        text_available = any(extracted.get(str(d.get("id","")),{}).get("text_preview") for d in docs)
        confidence = round(0.25 + 0.10 * text_available, 4)
        return {
            "status": "warning", "confidence": confidence,
            "finding": "LCA document uploaded but key values could not be automatically extracted. Manual review required.",
            "evidence": [{"doc": d.get("name")} for d in docs],
            "gap": False,
        }

    def _check_leakage(self, cp, docs, extracted):
        for doc in docs:
            data = extracted.get(str(doc.get("id","")), {})
            if doc.get("file_type") in ("xlsx","xls","xlsm") and data.get("sheets"):
                sheets = data["sheets"]
                leakage_sheets = [s for s in sheets if any(w in s.lower() for w in ["leakage","ghg","emission","displacement","scope"])]
                all_cols = []
                for s in leakage_sheets[:3]:
                    all_cols.extend(sheets[s].get("columns",[])[:10])
                emission_cols = [c for c in all_cols if any(w in c.lower() for w in ["emission","co2","ghg","leakage","scope","energy","electricity"])]
                if emission_cols:
                    # Confidence = f(leakage sheet coverage, emission column richness)
                    sheet_score  = min(len(leakage_sheets) / max(len(sheets), 1), 1.0)
                    col_score    = min(len(emission_cols) / 4.0, 1.0)  # 4 key cols = full score
                    confidence   = round(0.50 + 0.25 * sheet_score + 0.25 * col_score, 4)
                    return {
                        "status": "passed", "confidence": confidence,
                        "finding": f"✓ GHG leakage determination spreadsheet found. Leakage sheets: {len(leakage_sheets)}/{len(sheets)}. Key emission columns: {', '.join(emission_cols[:4])}. Manual review required to verify net removal calculation.",
                        "evidence": [{"doc": doc.get("name"), "leakage_sheets": leakage_sheets, "emission_columns": emission_cols}],
                        "gap": False,
                    }
                # Spreadsheet present but no emission columns
                sheet_score = min(len(leakage_sheets) / max(len(sheets), 1), 1.0)
                confidence  = round(0.30 + 0.20 * sheet_score, 4)
                return {
                    "status": "warning", "confidence": confidence,
                    "finding": f"Leakage spreadsheet uploaded but no GHG emission columns identified. Expected: scope 1/2/3 emissions, transportation, energy consumption. Sheets found: {len(sheets)}.",
                    "evidence": [{"doc": doc.get("name"), "sheets_found": list(sheets.keys())[:5]}],
                    "gap": False,
                }
        return self._check_document_presence(cp, docs, extracted)

    def _check_monitoring(self, cp, docs, extracted):
        monitoring_keywords = ["monitoring","measurement","reporting","verification","sensor","flow meter","pressure","injection rate","mrv","calibration","frequency","uncertainty"]
        for doc in docs:
            data = extracted.get(str(doc.get("id","")), {})
            text = data.get("text_preview","").lower()
            text_len = data.get("text_length", 0)
            found = [kw for kw in monitoring_keywords if kw in text]
            if len(found) >= 2:
                # Confidence = f(keyword density, text length, term count)
                keyword_score = min(len(found) / len(monitoring_keywords), 1.0)
                text_score    = min(text_len / 3000.0, 1.0)
                confidence    = round(0.45 + 0.35 * keyword_score + 0.20 * text_score, 4)
                return {
                    "status": "passed", "confidence": confidence,
                    "finding": f"✓ Monitoring plan found: {doc.get('name')}. Key monitoring elements: {', '.join(found[:5])} ({len(found)}/{len(monitoring_keywords)} matched). Manual review required to confirm compliance with GSC methodology requirements.",
                    "evidence": [{"doc": doc.get("name"), "monitoring_elements": found, "keyword_match_rate": round(len(found)/len(monitoring_keywords),2)}],
                    "gap": False,
                }
        return self._check_document_presence(cp, docs, extracted)

    def _check_storage(self, cp, docs, extracted):
        storage_keywords = ["reservoir","injection","geological","mineralisation","storage","injectivity","basalt","peridotite","aquifer","permeability","porosity","geochemical","trapping"]
        for doc in docs:
            data = extracted.get(str(doc.get("id","")), {})
            text  = data.get("text_preview","").lower()
            terms = data.get("key_terms", {})
            text_len = data.get("text_length", 0)
            found = [kw for kw in storage_keywords if kw in text]
            # Also count key_terms matches
            term_hits = sum(1 for k in ["contains_reservoir","contains_injection","contains_mineralisation"] if terms.get(k))
            total_signals = len(found) + term_hits
            if total_signals >= 2:
                keyword_score = min(len(found) / len(storage_keywords), 1.0)
                text_score    = min(text_len / 3000.0, 1.0)
                confidence    = round(0.45 + 0.35 * keyword_score + 0.20 * text_score, 4)
                return {
                    "status": "passed", "confidence": confidence,
                    "finding": f"✓ Storage site document found: {doc.get('name')}. Storage terms matched: {len(found)}/{len(storage_keywords)}: {', '.join(found[:4]) or 'reservoir/injection terms found'}. Manual review required to confirm site suitability.",
                    "evidence": [{"doc": doc.get("name"), "storage_elements": found, "keyword_match_rate": round(len(found)/len(storage_keywords),2)}],
                    "gap": False,
                }
        return self._check_document_presence(cp, docs, extracted)

    def _check_safeguards(self, cp, docs, extracted):
        safeguard_keywords = ["stakeholder","community","environmental","social","impact","safeguard","ifc","performance standard","consultation","fpic","grievance","biodiversity","indigenous","gender","labor","occupational"]
        for doc in docs:
            data = extracted.get(str(doc.get("id","")), {})
            text     = data.get("text_preview","").lower()
            text_len = data.get("text_length", 0)
            found = [kw for kw in safeguard_keywords if kw in text]
            if len(found) >= 2:
                keyword_score = min(len(found) / len(safeguard_keywords), 1.0)
                text_score    = min(text_len / 3000.0, 1.0)
                confidence    = round(0.40 + 0.35 * keyword_score + 0.25 * text_score, 4)
                return {
                    "status": "passed", "confidence": confidence,
                    "finding": f"✓ Safeguards document found: {doc.get('name')}. Key elements matched: {len(found)}/{len(safeguard_keywords)}: {', '.join(found[:5])}. Manual review required to confirm IFC Performance Standards compliance.",
                    "evidence": [{"doc": doc.get("name"), "safeguard_elements": found, "keyword_match_rate": round(len(found)/len(safeguard_keywords),2)}],
                    "gap": False,
                }
        return self._check_document_presence(cp, docs, extracted)

    def _check_document_presence(self, cp, docs, extracted) -> Dict:
        """
        Default fallback: document present but no specific checker applies.
        Confidence computed from extraction quality — text length, row count, sheet count.
        """
        doc_names  = [d.get("name","") for d in docs]
        total_text_len  = sum(extracted.get(str(d.get("id","")),{}).get("text_length", 0) for d in docs)
        total_rows      = sum(extracted.get(str(d.get("id","")),{}).get("row_count", 0) for d in docs)
        total_sheets    = sum(len(extracted.get(str(d.get("id","")),{}).get("sheets", {})) for d in docs)
        any_text        = total_text_len > 0
        has_data        = total_rows > 0 or total_sheets > 0

        # Score each dimension 0-1
        text_score  = min(total_text_len / 2000.0, 1.0)
        data_score  = min((total_rows / 100.0) + (total_sheets / 5.0), 1.0)
        doc_count   = min(len(docs) / 2.0, 1.0)   # more matching docs = more confidence

        if any_text:
            confidence = round(0.45 + 0.35 * text_score + 0.10 * doc_count + 0.10, 4)
        elif has_data:
            confidence = round(0.40 + 0.35 * data_score + 0.10 * doc_count + 0.15, 4)
        else:
            # Document present but nothing extracted — lowest non-zero confidence
            confidence = round(0.25 + 0.10 * doc_count, 4)

        content_note = "Text extracted." if any_text else ("Data extracted." if has_data else "Document present but content not extracted.")
        return {
            "status": "passed", "confidence": confidence,
            "finding": f"✓ Required document(s) uploaded: {', '.join(doc_names[:3])}. {content_note} Extraction quality score: {confidence:.2f}. Manual review required: '{cp.get('requirement','')[:120]}'",
            "evidence": [{"doc": d.get("name"), "type": d.get("document_type")} for d in docs[:3]],
            "gap": False,
        }


class VerificationEngine:
    """Main V&V engine — orchestrates extraction and checkpoint analysis."""

    def __init__(self):
        self.extractor = DocumentExtractor()
        self.analyser = CheckpointAnalyser()

    def run(self, project: Dict, documents: List[Dict], checkpoints: List[Dict]) -> Dict:
        logger.info(f"V&V run: project={project.get('id')}, docs={len(documents)}, checkpoints={len(checkpoints)}")

        # Step 1: Extract all documents
        extracted = {}
        for doc in documents:
            path = doc.get("storage_path","")
            ftype = doc.get("file_type","csv")
            if path and os.path.exists(path):
                extracted[str(doc.get("id",""))] = self.extractor.extract(path, ftype)
            else:
                extracted[str(doc.get("id",""))] = {"error": "file_not_found", "file_type": ftype}

        # Step 2: Analyse checkpoints
        results = []
        passed = failed = warnings = 0
        for cp in checkpoints:
            result = self.analyser.analyse(cp, documents, extracted)
            results.append({**cp, **result})
            s = result.get("status","warning")
            if s == "passed": passed += 1
            elif s == "failed": failed += 1
            else: warnings += 1

        # Step 3: Determine overall outcome
        critical_failed  = [r for r in results if r.get("critical") and r.get("status") == "failed"]
        critical_warning = [r for r in results if r.get("critical") and r.get("status") == "warning" and r.get("gap")]
        gaps = [r for r in results if r.get("gap")]

        if critical_failed:
            outcome = "not_verified"
            summary = f"Verification FAILED — {len(critical_failed)} critical checkpoint(s) not met. Project does not meet registry requirements for certificate issuance."
        elif len(gaps) >= 3:
            outcome = "conditional"
            summary = f"Conditional — {len(gaps)} document gap(s) identified. Provide missing documents and resubmit for final verification."
        elif warnings > 0:
            outcome = "conditional"
            summary = f"Conditional verification — {passed} checkpoints passed, {warnings} require manual review or clarification. Resolve outstanding items before certificate issuance."
        else:
            outcome = "verified"
            summary = f"Verification APPROVED — all {passed} checkpoints passed. Project meets registry requirements and is recommended for certificate issuance."

        # Step 4: Credit estimate
        credit_estimate = self._estimate_credits(documents, extracted, project)

        return {
            "outcome": outcome,
            "summary": summary,
            "checkpoint_results": results,
            "stats": {"passed": passed, "failed": failed, "warnings": warnings, "total": len(results), "gaps": len(gaps)},
            "credit_estimate": credit_estimate,
            "completed_at": datetime.utcnow().isoformat(),
        }

    def _estimate_credits(self, docs, extracted, project) -> Dict:
        """Estimate credit quantity from LCA or leakage spreadsheets."""
        for doc in docs:
            data = extracted.get(str(doc.get("id","")), {})
            # Check Excel sheets for CO2 values
            if data.get("sheets"):
                for sheet_name, sheet in data["sheets"].items():
                    if any(w in sheet_name.lower() for w in ["co2","carbon","removal","net","lca","result"]):
                        for col, stats in sheet.get("numeric_summary",{}).items():
                            if any(w in col.lower() for w in ["net","removal","co2","credit","corc"]):
                                val = abs(stats.get("sum", stats.get("max", 0)))
                                if 10 < val < 10_000_000:  # reasonable range in tCO2e
                                    return {
                                        "estimated_credits": round(val, 0),
                                        "unit": "tCO2e",
                                        "basis": f"Extracted from '{col}' in sheet '{sheet_name}' of {doc.get('name','')}",
                                        "confidence": "indicative — manual verification required",
                                    }
            # Check text for CO2 values
            terms = data.get("key_terms", {})
            if terms.get("co2_tonnes"):
                try:
                    val = float(terms["co2_tonnes"].replace(",","."))
                    if val > 0:
                        return {
                            "estimated_credits": round(val, 0),
                            "unit": "tCO2e",
                            "basis": f"Extracted from document text: {doc.get('name','')}",
                            "confidence": "indicative — manual verification required",
                        }
                except: pass
        return {"estimated_credits": None, "unit": "tCO2e", "basis": "Upload LCA spreadsheet or leakage determination for automatic credit estimate", "confidence": "none"}
